#!/usr/bin/env python3
"""
Office文件表格提取工具 - 从.docx, .pptx, .xlsx文件中提取所有表格为CSV
版本: 2.0.0
作者: tianli
"""

import sys
import argparse
import pandas as pd
from pathlib import Path
from docx import Document
from pptx import Presentation

from common_utils import (
    show_success, show_error, show_warning, show_info, show_processing,
    validate_input_file, ensure_directory, ProgressTracker, fatal_error,
    check_python_packages, find_files_by_extension
)

SCRIPT_VERSION = "2.0.0"

def check_dependencies():
    show_info("检查依赖项...")
    if not check_python_packages(['pandas', 'python-docx', 'python-pptx', 'openpyxl']):
        sys.exit(1)
    show_success("依赖检查完成")

def extract_from_docx(file_path: Path, output_dir: Path) -> int:
    try:
        doc = Document(file_path)
        if not doc.tables:
            return 0
        
        count = 0
        for i, table in enumerate(doc.tables, 1):
            data = [[cell.text for cell in row.cells] for row in table.rows]
            df = pd.DataFrame(data)
            output_file = output_dir / f"{file_path.stem}_table_{i}.csv"
            df.to_csv(output_file, index=False, header=False)
            count += 1
        return count
    except Exception as e:
        show_error(f"处理DOCX失败: {e}")
        return 0

def extract_from_pptx(file_path: Path, output_dir: Path) -> int:
    try:
        prs = Presentation(file_path)
        count = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                
                table = shape.table
                data = [[cell.text for cell in row.cells] for row in table.rows]
                df = pd.DataFrame(data)
                output_file = output_dir / f"{file_path.stem}_slide_{slide_num}_table_{count+1}.csv"
                df.to_csv(output_file, index=False, header=False)
                count += 1
        return count
    except Exception as e:
        show_error(f"处理PPTX失败: {e}")
        return 0

def extract_from_xlsx(file_path: Path, output_dir: Path) -> int:
    try:
        xls = pd.ExcelFile(file_path)
        count = 0
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
            output_file = output_dir / f"{file_path.stem}_sheet_{sheet_name}.csv"
            df.to_csv(output_file, index=False, header=False)
            count += 1
        return count
    except Exception as e:
        show_error(f"处理XLSX失败: {e}")
        return 0

def extract_tables_from_file(file_path: Path, output_dir: Path) -> int:
    if not validate_input_file(file_path):
        return 0

    ext = file_path.suffix.lower()
    extractors = {
        '.docx': extract_from_docx,
        '.pptx': extract_from_pptx,
        '.xlsx': extract_from_xlsx
    }

    if ext not in extractors:
        show_warning(f"跳过不支持的文件类型: {file_path.name}")
        return 0
        
    file_output_dir = output_dir / file_path.stem
    ensure_directory(file_output_dir)
    
    show_processing(f"从 {file_path.name} 提取表格...")
    count = extractors[ext](file_path, file_output_dir)
    
    if count > 0:
        show_success(f"成功提取 {count} 个表格到 {file_output_dir}")
    else:
        show_info(f"在 {file_path.name} 中未找到表格")
        
    return count

def main():
    parser = argparse.ArgumentParser(description="Office文件表格提取工具")
    parser.add_argument("input_paths", nargs='+', help="一个或多个文件/目录路径")
    parser.add_argument("-o", "--output", help="输出目录 (默认: ./extracted_tables)")
    parser.add_argument("-r", "--recursive", action="store_true", help="递归处理目录")
    parser.add_argument('--version', action='version', version=f'%(prog)s {SCRIPT_VERSION}')
    args = parser.parse_args()

    check_dependencies()

    output_dir = Path(args.output) if args.output else Path("./extracted_tables")
    ensure_directory(output_dir)
    show_info(f"输出目录: {output_dir.resolve()}")

    files_to_process = find_files_by_extension(
        args.input_paths,
        ['docx', 'pptx', 'xlsx'],
        recursive=args.recursive
    )

    if not files_to_process:
        show_warning("未找到任何支持的Office文件")
        sys.exit(0)

    total_extracted = 0
    progress = ProgressTracker(len(files_to_process))

    for file_path in files_to_process:
        progress.show(f"处理 {file_path.name}")
        count = extract_tables_from_file(file_path, output_dir)
        total_extracted += count
        
    show_info("\n处理完成")
    show_success(f"总共提取了 {total_extracted} 个表格")

if __name__ == "__main__":
    main()

