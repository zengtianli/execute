#!/usr/bin/env python3
"""
PPTX转Markdown转换工具 - 将PPTX演示文稿转换为Markdown格式
版本: 2.0.0
作者: tianli
"""

import sys
import argparse
from pathlib import Path
from pptx import Presentation

from common_utils import (
    show_success, show_error, show_warning, show_info, show_processing,
    validate_input_file, ensure_directory, ProgressTracker, fatal_error,
    check_python_packages, find_files_by_extension, get_file_basename
)

SCRIPT_VERSION = "2.0.0"

def check_dependencies():
    show_info("检查依赖项...")
    if not check_python_packages(['python-pptx']):
        sys.exit(1)
    show_success("依赖检查完成")

def convert_pptx_to_md_single(file_path: Path, output_dir: Path) -> bool:
    if not validate_input_file(file_path):
        return False

    if file_path.suffix.lower() != '.pptx':
        show_warning(f"跳过非PPTX文件: {file_path.name}")
        return False

    base_name = get_file_basename(file_path)
    md_output_dir = output_dir / base_name
    ensure_directory(md_output_dir)
    output_file = md_output_dir / f"{base_name}.md"
    
    show_processing(f"转换 {file_path.name} 为 Markdown...")

    try:
        prs = Presentation(file_path)
        with open(output_file, 'w', encoding='utf-8') as md_file:
            for i, slide in enumerate(prs.slides, 1):
                md_file.write(f"## Slide {i}\n\n")
                
                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        md_file.write(shape.text + '\n\n')

                if notes:
                    md_file.write(f"### Speaker Notes\n\n{notes}\n\n")
                
                md_file.write("---\n\n")
        
        show_success(f"成功转换: {file_path.name} -> {output_file}")
        return True
    except Exception as e:
        show_error(f"转换失败 {file_path.name}: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(description="PPTX转Markdown转换工具")
    parser.add_argument("input_paths", nargs='+', help="一个或多个PPTX文件/目录路径")
    parser.add_argument("-o", "--output", help="输出目录 (默认: ./converted_md)")
    parser.add_argument("-r", "--recursive", action="store_true", help="递归处理目录")
    parser.add_argument('--version', action='version', version=f'%(prog)s {SCRIPT_VERSION}')
    args = parser.parse_args()

    check_dependencies()

    output_dir = Path(args.output) if args.output else Path("./converted_md")
    ensure_directory(output_dir)
    show_info(f"输出目录: {output_dir.resolve()}")

    files_to_process = find_files_by_extension(
        args.input_paths,
        ['pptx'],
        recursive=args.recursive
    )

    if not files_to_process:
        show_warning("未找到任何PPTX文件")
        sys.exit(0)

    total_success = 0
    progress = ProgressTracker(len(files_to_process))

    for file_path in files_to_process:
        progress.show(f"处理 {file_path.name}")
        if convert_pptx_to_md_single(file_path, output_dir):
            total_success += 1
            
    show_info("\n处理完成")
    show_success(f"总共成功转换了 {total_success} 个文件")

if __name__ == "__main__":
    main()

